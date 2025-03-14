import os
import subprocess
import uuid
import io
import zipfile
import threading
import time
from flask import Flask, render_template, request, send_file, jsonify, abort
from pdf2image import convert_from_path
from pptx import Presentation
from werkzeug.utils import secure_filename  # For sanitizing filenames

app = Flask(__name__)

# Set folders for uploads and converted files.
UPLOAD_FOLDER = "uploads"
CONVERTED_FOLDER = "converted_files"

# Ensure necessary folders exist.
for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER]:
    if not os.path.exists(folder):
        os.makedirs(folder)

# Global dictionaries to track file metadata.
# uploads_tracking: maps the stored filename to a dict with upload time and the original name.
uploads_tracking = {}  # Format: {stored_filename: {"upload_time": timestamp, "original_name": original_filename}}
# converted_tracking: maps the converted file's name to its last access time.
converted_tracking = {}  # Format: {filename: last_access_time}

# Expiration thresholds in seconds.
UPLOAD_EXPIRY = 60         # 1 hour for files in the uploads folder.
CONVERTED_EXPIRY = 300       # 5 minutes for files in the converted folder.

def generate_nicer_filename(original_name, new_extension):
    """
    Generates a nicer filename using a sanitized version of the original name and a timestamp.
    For example, "My File.pdf" with new_extension ".pptx" may become "My_File_1678901234.pptx".
    """
    # Sanitize the original filename.
    base = secure_filename(original_name)
    # Remove the original extension.
    base = os.path.splitext(base)[0]
    # Append current timestamp to help ensure uniqueness.
    timestamp = int(time.time())
    return f"{base}_{timestamp}{new_extension}"

def cleanup_files():
    """
    Background thread that periodically checks for files in UPLOAD_FOLDER and CONVERTED_FOLDER
    that have exceeded their allowed lifetime and removes them.
    """
    while True:
        current_time = time.time()

        # Check for expired files in UPLOAD_FOLDER.
        expired_uploads = [fname for fname, meta in uploads_tracking.items()
                           if current_time - meta["upload_time"] > UPLOAD_EXPIRY]
        for fname in expired_uploads:
            file_path = os.path.join(UPLOAD_FOLDER, fname)
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    print(f"Removed expired uploaded file: {fname}")
                uploads_tracking.pop(fname, None)
            except Exception as e:
                print(f"Error removing uploaded file {fname}: {e}")

        # Check for expired files in CONVERTED_FOLDER.
        expired_converted = [fname for fname, last_access in converted_tracking.items()
                             if current_time - last_access > CONVERTED_EXPIRY]
        for fname in expired_converted:
            file_path = os.path.join(CONVERTED_FOLDER, fname)
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    print(f"Removed expired converted file: {fname}")
                converted_tracking.pop(fname, None)
            except Exception as e:
                print(f"Error removing converted file {fname}: {e}")

        time.sleep(30)

# Start the background cleanup thread.
cleanup_thread = threading.Thread(target=cleanup_files, daemon=True)
cleanup_thread.start()

def save_uploaded_file(file_obj, ext):
    """
    Saves the uploaded file in UPLOAD_FOLDER.
    Instead of using a plain UUID, the file is stored with a unique name,
    while also storing the original name for a nicer naming during conversion.
    """
    # Generate a unique name for internal storage.
    unique_name = f"{uuid.uuid4()}{ext}"
    file_path = os.path.join(UPLOAD_FOLDER, unique_name)
    file_obj.save(file_path)
    # Store metadata including the original filename.
    uploads_tracking[unique_name] = {"upload_time": time.time(), "original_name": file_obj.filename}
    return file_path

def convert_pdf_to_pptx(pdf_path, original_name=None):
    """
    Converts a PDF to PPTX.
    If an original name is provided, a nicer filename is generated for the PPTX.
    Otherwise, a UUID-based name is used.
    """
    images = convert_from_path(pdf_path)
    presentation = Presentation()
    for image in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        iw, ih = image.size
        sw, sh = presentation.slide_width, presentation.slide_height
        ratio = iw / ih
        if ratio > (sw / sh):
            nw = sw
            nh = sw / ratio
        else:
            nh = sh
            nw = sh * ratio
        left = (sw - nw) / 2
        top = (sh - nh) / 2
        img_bytes = io.BytesIO()
        image.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        slide.shapes.add_picture(img_bytes, left, top, width=nw, height=nh)
    # Generate a nicer filename if original_name is available.
    if original_name:
        nice_filename = generate_nicer_filename(original_name, ".pptx")
    else:
        nice_filename = f"{uuid.uuid4()}.pptx"
    output_path = os.path.join(CONVERTED_FOLDER, nice_filename)
    presentation.save(output_path)
    converted_tracking[nice_filename] = time.time()
    return nice_filename

def convert_pptx_to_pdf(pptx_path, original_name=None):
    """
    Converts a PPTX file to PDF using LibreOffice in headless mode.
    Generates a nicer filename based on the original name if available.
    """
    output_folder = os.path.abspath(CONVERTED_FOLDER)
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", output_folder],
        check=True
    )
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    output_filename = f"{base}.pdf"
    if original_name:
        nice_filename = generate_nicer_filename(original_name, ".pdf")
    else:
        nice_filename = f"{uuid.uuid4()}.pdf"
    src = os.path.join(output_folder, output_filename)
    dst = os.path.join(output_folder, nice_filename)
    os.rename(src, dst)
    converted_tracking[nice_filename] = time.time()
    return nice_filename

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/result")
def result():
    return render_template("result.html")

@app.route("/upload", methods=["POST"])
def upload():
    """
    Handles file uploads. Verifies file size and conversion type,
    then stores the file while keeping track of its original name.
    """
    file = request.files.get("file")
    conversion_type = request.form.get("conversion_type")
    if not file:
        return jsonify({"error": "No file uploaded"}), 400
    if file.content_length is not None and file.content_length > 104857600:  # 100 MB limit.
        return jsonify({"error": "File size exceeds 100MB"}), 400
    if conversion_type == "pdf_to_pptx":
        if not file.filename.lower().endswith(".pdf"):
            return jsonify({"error": "Only PDF files allowed"}), 400
    elif conversion_type == "pptx_to_pdf":
        if not file.filename.lower().endswith(".pptx"):
            return jsonify({"error": "Only PPTX files allowed"}), 400
    else:
        return jsonify({"error": "Invalid conversion type"}), 400
    saved_path = save_uploaded_file(file, os.path.splitext(file.filename)[1])
    return jsonify({"file_id": os.path.basename(saved_path), "original_name": file.filename})

@app.route("/convert_all", methods=["POST"])
def convert_all():
    """
    Processes the list of file IDs for conversion.
    Uses the original filename from the metadata to generate nicer names for the converted files.
    Returns a direct download link for a single file or a ZIP archive for multiple files.
    """
    data = request.get_json()
    conversion_type = data.get("conversion_type")
    file_ids = data.get("file_ids")
    if not file_ids or not isinstance(file_ids, list):
        return jsonify({"error": "No files provided"}), 400
    converted_files = []
    try:
        if conversion_type == "pdf_to_pptx":
            for fid in file_ids:
                pdf_path = os.path.join(UPLOAD_FOLDER, fid)
                if not os.path.exists(pdf_path):
                    continue
                # Retrieve original name from tracking.
                original_name = uploads_tracking.get(fid, {}).get("original_name")
                conv_filename = convert_pdf_to_pptx(pdf_path, original_name)
                converted_files.append(conv_filename)
                uploads_tracking.pop(fid, None)
                os.remove(pdf_path)
        elif conversion_type == "pptx_to_pdf":
            for fid in file_ids:
                pptx_path = os.path.join(UPLOAD_FOLDER, fid)
                if not os.path.exists(pptx_path):
                    continue
                original_name = uploads_tracking.get(fid, {}).get("original_name")
                conv_filename = convert_pptx_to_pdf(pptx_path, original_name)
                converted_files.append(conv_filename)
                uploads_tracking.pop(fid, None)
                os.remove(pptx_path)
        else:
            return jsonify({"error": "Invalid conversion type"}), 400

        if not converted_files:
            return jsonify({"error": "No files converted"}), 400

        # For multiple files, generate a nicer ZIP filename.
        if len(converted_files) == 1:
            download_url = f"/download/{converted_files[0]}"
            return jsonify({"download_url": download_url})
        else:
            # Use a nicer name for the ZIP archive.
            nice_zip_name = generate_nicer_filename("converted_files", ".zip")
            zip_path = os.path.join(CONVERTED_FOLDER, nice_zip_name)
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for conv in converted_files:
                    conv_path = os.path.join(CONVERTED_FOLDER, conv)
                    zipf.write(conv_path, arcname=conv)
                    os.remove(conv_path)
                    converted_tracking.pop(conv, None)
            converted_tracking[nice_zip_name] = time.time()
            download_url = f"/download/{nice_zip_name}"
            return jsonify({"download_url": download_url})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/download/<filename>")
def download(filename):
    """
    Serves files from the CONVERTED_FOLDER.
    Instead of immediate removal, updates the last access time so that the background
    cleanup process handles deletion safely.
    """
    file_path = os.path.join(CONVERTED_FOLDER, filename)
    if not os.path.exists(file_path):
        abort(404)
    converted_tracking[filename] = time.time()
    return send_file(file_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)